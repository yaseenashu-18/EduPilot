import os
import sys
import subprocess

# 1. Programmatic installation of python-pptx if not present
try:
    import pptx
except ImportError:
    print("python-pptx package not found. Installing now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 2. Initialize presentation with widescreen 16:9 format
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 3. Design Tokens (Premium Dark Slate Theme)
    BG_COLOR = RGBColor(11, 15, 25)          # #0B0F19 - Deep Space Navy
    TEXT_LIGHT = RGBColor(248, 250, 252)    # #F8FAFC - Ice White
    TEXT_MUTED = RGBColor(148, 163, 184)    # #94A3B8 - Slate Muted
    ACCENT_SKY = RGBColor(14, 165, 233)     # #0EA5E9 - Sky Blue
    ACCENT_INDIGO = RGBColor(79, 70, 229)   # #4F46E5 - Premium Indigo
    ACCENT_RED = RGBColor(239, 68, 68)       # #EF4444 - Warning Red
    ACCENT_GREEN = RGBColor(34, 197, 94)     # #22C55E - Success Green
    CARD_BG = RGBColor(21, 29, 45)           # #151D2D - Glassmorphic Card Fill
    CARD_BORDER = RGBColor(38, 48, 71)       # #263047 - Subtle Card Border
    
    # Clean blank layout
    blank_layout = prs.slide_layouts[6]
    
    def apply_solid_background(slide, color):
        """Applies a full-screen solid colored rectangle as the background"""
        bg_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = color
        bg_rect.line.fill.background() # No border
        return bg_rect

    def add_slide_header(slide, title_text, category_text="EDUPILOT AI"):
        """Adds a standardized premium header to content slides"""
        # Category label (small caps style)
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Segoe UI"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_SKY
        
        # Main slide title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT

    def create_glass_card(slide, left, top, width, height, border_color=CARD_BORDER, fill_color=CARD_BG):
        """Creates a glassmorphic card shape to host structured data"""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide1, BG_COLOR)
    
    # Dynamic decorative glowing orbits
    decor = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    decor.fill.solid()
    decor.fill.fore_color.rgb = RGBColor(20, 24, 45)
    decor.line.fill.background()
    
    decor2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), Inches(6), Inches(6))
    decor2.fill.solid()
    decor2.fill.fore_color.rgb = RGBColor(18, 30, 48)
    decor2.line.fill.background()
    
    # Title Text Box
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "EDUPILOT AI"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    
    # Accent indicator line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.2), Inches(2.0), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_SKY
    line.line.fill.background()
    
    # Subtitle Text Box
    sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(11.0), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p2 = tf_sub.paragraphs[0]
    p2.text = "Intellectual Learning & Performance Analytics\nPowered by Machine Learning"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_SKY
    
    # Presentation Footer details
    footer_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(11.0), Inches(0.8))
    tf_foot = footer_box.text_frame
    tf_foot.word_wrap = True
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = "PARUL UNIVERSITY CAMPUS INTERVENTION PLATFORM\nFull-Stack Machine Learning Solution (Python Flask / MongoDB / Scikit-Learn)"
    p_foot.font.name = "Segoe UI"
    p_foot.font.size = Pt(11)
    p_foot.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: THE CHALLENGE
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide2, BG_COLOR)
    add_slide_header(slide2, "The Critical Gaps in Academic Operations", "Problem Statement")
    
    # Left Card: Student Issues
    create_glass_card(slide2, Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_RED)
    box_l = slide2.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Student Vulnerabilities"
    p_lh.font.name = "Segoe UI"
    p_lh.font.size = Pt(20)
    p_lh.font.bold = True
    p_lh.font.color.rgb = ACCENT_RED
    p_lh.space_after = Pt(14)
    
    bullets_l = [
        "Undetected Academic Risk: Students fail to notice sliding averages or attendance shortages until it bars them from examinations.",
        "Generic Goal Alignment: Standard curriculums treat all students identically, offering no automated guidance for ambitious job roles.",
        "Skill Disconnect: Ambitious career aspirations (e.g. AIML, Cyber) are targeted without early, quantitative skill metrics analysis."
    ]
    for bullet in bullets_l:
        p = tf_l.add_paragraph()
        p.text = "• " + bullet
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)
        
    # Right Card: Institutional Gaps
    create_glass_card(slide2, Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_INDIGO)
    box_r = slide2.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Institutional Hurdles"
    p_rh.font.name = "Segoe UI"
    p_rh.font.size = Pt(20)
    p_rh.font.bold = True
    p_rh.font.color.rgb = ACCENT_INDIGO
    p_rh.space_after = Pt(14)
    
    bullets_r = [
        "Supervision Deficits: Professors manage hundreds of students, making personalized proactive mentoring mathematically impossible.",
        "Disconnected Data Repositories: Academics, programming milestones, and class attendances sit in manual, unlinked spreadsheets.",
        "Reactive Decision Models: Institutions review failure rates retrospectively instead of predicting risk dynamically in real-time."
    ]
    for bullet in bullets_r:
        p = tf_r.add_paragraph()
        p.text = "• " + bullet
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 3: THE SOLUTION
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide3, BG_COLOR)
    add_slide_header(slide3, "EduPilot AI Solution Architecture", "Core Ecosystem")
    
    # 4 Solutions Cards Grid
    grid_coords = [
        (Inches(1.0), Inches(1.8), "🔮 Automated Predictions", "Employs Random Forest and Linear Regression models to evaluate student risk levels (High, Medium, Low) and forecast future GPAs.", ACCENT_SKY),
        (Inches(7.0), Inches(1.8), "🎯 Dynamic Weekly Tasks", "Recommends weak-subject-oriented checklists using a rule engine that updates status and registers progress in real-time.", ACCENT_INDIGO),
        (Inches(1.0), Inches(4.4), "📊 Competency Profiling", "Renders rich interactive Chart.js visualizations outlining marks, GPA trends, and radar competence views.", ACCENT_GREEN),
        (Inches(7.0), Inches(4.4), "🔑 Institutional Control", "Hosts a dedicated Administrative Command Center to retrain ML models and manage secure faculty accounts dynamically.", ACCENT_SKY)
    ]
    
    for x, y, title, desc, border_col in grid_coords:
        create_glass_card(slide3, x, y, Inches(5.3), Inches(2.2), border_color=border_col)
        box = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(4.9), Inches(1.9))
        tf_box = box.text_frame
        tf_box.word_wrap = True
        
        p_h = tf_box.paragraphs[0]
        p_h.text = title
        p_h.font.name = "Segoe UI"
        p_h.font.size = Pt(18)
        p_h.font.bold = True
        p_h.font.color.rgb = border_col
        p_h.space_after = Pt(6)
        
        p_d = tf_box.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Segoe UI"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_LIGHT

    # =========================================================================
    # SLIDE 4: SYSTEM TOPOLOGY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide4, BG_COLOR)
    add_slide_header(slide4, "Full-Stack Deployment & Database Model", "System Topology")
    
    # Left column: Tech Stack Card
    create_glass_card(slide4, Inches(1.0), Inches(1.8), Inches(4.5), Inches(4.8), border_color=ACCENT_SKY)
    box_ts = slide4.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.1), Inches(4.4))
    tf_ts = box_ts.text_frame
    tf_ts.word_wrap = True
    
    p_tsh = tf_ts.paragraphs[0]
    p_tsh.text = "Software Stack & Integration"
    p_tsh.font.name = "Segoe UI"
    p_tsh.font.size = Pt(18)
    p_tsh.font.bold = True
    p_tsh.font.color.rgb = ACCENT_SKY
    p_tsh.space_after = Pt(14)
    
    techs = [
        "Backend Framework: Python Flask REST Controller",
        "Frontend Layout: Modern HTML5, Vanilla CSS3 (Glassmorphic variables), Bootstrap 5",
        "Database Client: PyMongo interfacing with a cloud MongoDB cluster (with SQLite fallback setup)",
        "Identity Layer: Google OAuth 2.0 Identity Flow with automated campus-restricted domain gateway checking",
        "Mailer Gateway: Robust secure SMTP protocol for automated token-based password resets"
    ]
    for tech in techs:
        p = tf_ts.add_paragraph()
        p.text = "✔  " + tech
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)
        
    # Right Columns: Mongo Collections Representation
    create_glass_card(slide4, Inches(6.0), Inches(1.8), Inches(6.3), Inches(4.8), border_color=ACCENT_INDIGO)
    box_db = slide4.shapes.add_textbox(Inches(6.2), Inches(2.0), Inches(5.9), Inches(4.4))
    tf_db = box_db.text_frame
    tf_db.word_wrap = True
    
    p_dbh = tf_db.paragraphs[0]
    p_dbh.text = "Centralized MongoDB Collections Schema"
    p_dbh.font.name = "Segoe UI"
    p_dbh.font.size = Pt(18)
    p_dbh.font.bold = True
    p_dbh.font.color.rgb = ACCENT_INDIGO
    p_dbh.space_after = Pt(14)
    
    collections = [
        ("students", "Stores secure hashed passwords, active demographics, cumulative CGPA, attendance logs, and selected career targets."),
        ("teachers", "Registers faculty accounts, designated supervision departments, and profiles authorized to review class rosters."),
        ("subjects", "Documents active academic courses per student alongside raw internal scoring metrics (out of 100)."),
        ("tasks", "Tracks individualized weekly checklists containing task descriptions, strict target timelines, and completion flags."),
        ("predictions", "Stores pre-calculated placement readiness ratios, anticipated GPAs, and risk levels calculated via Scikit-Learn.")
    ]
    
    for coll_name, coll_desc in collections:
        p = tf_db.add_paragraph()
        p.text = f"📂 db.{coll_name} → "
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_SKY
        
        # append description in same paragraph
        run = p.add_run()
        run.text = coll_desc
        run.font.name = "Segoe UI"
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 5: STUDENT COCKPIT
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide5, BG_COLOR)
    add_slide_header(slide5, "The Student Console – Performance Cockpit", "User Experience (UX)")
    
    # Visual grid of cards
    create_glass_card(slide5, Inches(1.0), Inches(1.8), Inches(3.4), Inches(2.2), border_color=ACCENT_SKY)
    box1 = slide5.shapes.add_textbox(Inches(1.15), Inches(1.95), Inches(3.1), Inches(1.9))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "📈 Unified Roster & KPIs"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_SKY
    p1.space_after = Pt(6)
    p1_d = tf1.add_paragraph()
    p1_d.text = "Monitors real-time attendance thresholds, cumulative CGPA, assignment progress, and coding streak indicators simultaneously."
    p1_d.font.name = "Segoe UI"
    p1_d.font.size = Pt(12)
    p1_d.font.color.rgb = TEXT_LIGHT
    
    create_glass_card(slide5, Inches(4.8), Inches(1.8), Inches(3.4), Inches(2.2), border_color=ACCENT_INDIGO)
    box2 = slide5.shapes.add_textbox(Inches(4.95), Inches(1.95), Inches(3.1), Inches(1.9))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "🎯 Smart Weak-Subject Checklists"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_INDIGO
    p2.space_after = Pt(6)
    p2_d = tf2.add_paragraph()
    p2_d.text = "Actionable tasks are automatically customized for subjects where scoring lies below 60%. Checking items updates database states."
    p2_d.font.name = "Segoe UI"
    p2_d.font.size = Pt(12)
    p2_d.font.color.rgb = TEXT_LIGHT
    
    create_glass_card(slide5, Inches(8.6), Inches(1.8), Inches(3.7), Inches(2.2), border_color=ACCENT_GREEN)
    box3 = slide5.shapes.add_textbox(Inches(8.75), Inches(1.95), Inches(3.4), Inches(1.9))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "🔮 Interactive Charts"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN
    p3.space_after = Pt(6)
    p3_d = tf3.add_paragraph()
    p3_d.text = "Generates responsive, high-fidelity Chart.js visualizations outlining grade trajectory trends and multi-dimensional radar plots."
    p3_d.font.name = "Segoe UI"
    p3_d.font.size = Pt(12)
    p3_d.font.color.rgb = TEXT_LIGHT
    
    # Large bottom dashboard detail text block
    create_glass_card(slide5, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.3), border_color=CARD_BORDER)
    box_b = slide5.shapes.add_textbox(Inches(1.2), Inches(4.45), Inches(10.9), Inches(2.0))
    tf_b = box_b.text_frame
    tf_b.word_wrap = True
    p_bh = tf_b.paragraphs[0]
    p_bh.text = "Adaptive Career Alignment"
    p_bh.font.name = "Segoe UI"
    p_bh.font.size = Pt(16)
    p_bh.font.bold = True
    p_bh.font.color.rgb = ACCENT_SKY
    p_bh.space_after = Pt(6)
    
    p_bd = tf_b.add_paragraph()
    p_bd.text = "Students can dynamically select their ultimate professional objectives from 5 distinct disciplines:\n" \
                "  • AIML Engineer: Challenges them with Python data structures and PyTorch pipeline projects.\n" \
                "  • Software Developer: Recommends OOP paradigms, version controls, and structured LeetCode problems.\n" \
                "  • Data Analyst: Maps SQL aggregate querying tasks alongside data storytelling benchmarks.\n" \
                "  • DevOps & Cybersecurity: Directs Linux administration scripting, containerization (Docker/Kubernetes), and OWASP lab practices."
    p_bd.font.name = "Segoe UI"
    p_bd.font.size = Pt(12)
    p_bd.font.color.rgb = TEXT_LIGHT
    p_bd.space_after = Pt(4)

    # =========================================================================
    # SLIDE 6: THE FACULTY PANEL
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide6, BG_COLOR)
    add_slide_header(slide6, "The Faculty Panel – Supervision & Risk Indicators", "User Experience (UX)")
    
    create_glass_card(slide6, Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_SKY)
    box_fl = slide6.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_fl = box_fl.text_frame
    tf_fl.word_wrap = True
    
    p_flh = tf_fl.paragraphs[0]
    p_flh.text = "Supervision Interface & Filters"
    p_flh.font.name = "Segoe UI"
    p_flh.font.size = Pt(18)
    p_flh.font.bold = True
    p_flh.font.color.rgb = ACCENT_SKY
    p_flh.space_after = Pt(14)
    
    fl_points = [
        "Dynamic Supervision Registry: Renders a comprehensive list of all students enrolled within the faculty's specialized department.",
        "Smart Search & Filter System: Enables search by student name or immediate performance sorting based on current CGPA tiers.",
        "Classroom Aggregate Insight: Features high-level counter displays tracking total supervised students, average class CGPA, and specific attendance thresholds."
    ]
    for pt in fl_points:
        p = tf_fl.add_paragraph()
        p.text = "✔  " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)
        
    create_glass_card(slide6, Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_INDIGO)
    box_fr = slide6.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_fr = box_fr.text_frame
    tf_fr.word_wrap = True
    
    p_frh = tf_fr.paragraphs[0]
    p_frh.text = "Proactive Risk Alert System"
    p_frh.font.name = "Segoe UI"
    p_frh.font.size = Pt(18)
    p_frh.font.bold = True
    p_frh.font.color.rgb = ACCENT_INDIGO
    p_frh.space_after = Pt(14)
    
    fr_points = [
        "ML-Driven Risk Labels: Instantly classifies and color-codes students based on real-time parameter assessments (High, Medium, Low risk badges).",
        "Direct Intervention Triggers: Provides immediate visual prompts flagging critical status alerts (e.g. attendance shortages under 70%).",
        "PDF Performance Exporter: Features clean layout report templates allowing professors to trigger structured PDF print sheets containing academic metrics for offline review."
    ]
    for pt in fr_points:
        p = tf_fr.add_paragraph()
        p.text = "✔  " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 7: ADMINISTRATIVE CONTROL
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide7, BG_COLOR)
    add_slide_header(slide7, "Administrative Console & ML Pipeline Retraining", "Institutional Control")
    
    # Left Box: Global Registry Configuration
    create_glass_card(slide7, Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_INDIGO)
    box_al = slide7.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_al = box_al.text_frame
    tf_al.word_wrap = True
    
    p_alh = tf_al.paragraphs[0]
    p_alh.text = "Global Registry Configuration"
    p_alh.font.name = "Segoe UI"
    p_alh.font.size = Pt(18)
    p_alh.font.bold = True
    p_alh.font.color.rgb = ACCENT_INDIGO
    p_alh.space_after = Pt(14)
    
    al_points = [
        "Faculty Credential Management: Administrators host a dedicated control panel to Add, Edit, or Terminate professor profiles.",
        "Strict Domain Compliance: Gateway verification ensuring all created profiles match required corporate university standards.",
        "System Telemetry Dashboard: Tracks total registered students, active professors, and global telemetry profiles in real-time."
    ]
    for pt in al_points:
        p = tf_al.add_paragraph()
        p.text = "✔  " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)
        
    # Right Box: Machine Learning Lifecycle Management
    create_glass_card(slide7, Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_SKY)
    box_ar = slide7.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_ar = box_ar.text_frame
    tf_ar.word_wrap = True
    
    p_arh = tf_ar.paragraphs[0]
    p_arh.text = "Machine Learning Lifecycle Control"
    p_arh.font.name = "Segoe UI"
    p_arh.font.size = Pt(18)
    p_arh.font.bold = True
    p_arh.font.color.rgb = ACCENT_SKY
    p_arh.space_after = Pt(14)
    
    ar_points = [
        "1. Dynamic Record Generation: Re-triggers dataset pipelines generating 5,000 synthetic student entries containing highly correlated indicators (NumPy & Faker).",
        "2. Automated Pipeline Retraining: Instantly trains Random Forest classifiers and Linear Regression engines on fresh datasets.",
        "3. Real-Time Serialization: Saves trained parameters directly to backend directories as serialized Joblib files (.pkl) instantly ready for active inferences."
    ]
    for pt in ar_points:
        p = tf_ar.add_paragraph()
        p.text = pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 8: MACHINE LEARNING ENGINE
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide8, BG_COLOR)
    add_slide_header(slide8, "Under the Hood – Machine Learning Algorithms", "Machine Learning Core")
    
    # Left Model: Random Forest
    create_glass_card(slide8, Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_SKY)
    box_ml1 = slide8.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_ml1 = box_ml1.text_frame
    tf_ml1.word_wrap = True
    
    p_ml1h = tf_ml1.paragraphs[0]
    p_ml1h.text = "Academic Risk Classification"
    p_ml1h.font.name = "Segoe UI"
    p_ml1h.font.size = Pt(18)
    p_ml1h.font.bold = True
    p_ml1h.font.color.rgb = ACCENT_SKY
    p_ml1h.space_after = Pt(14)
    
    points_ml1 = [
        "Algorithm: RandomForestClassifier (Scikit-Learn).",
        "Input Indicators (4-dimensional): Attendance, Cumulative CGPA, Coding Streak score, and Assignment progress.",
        "Hyperparameters: 100 estimators, max depth of 10, stratified train-test splits ensuring class balance.",
        "Output Metrics: Assigns definitive labels ('High Risk', 'Medium Risk', or 'Low Risk') to active profiles.",
        "Robust Fallbacks: Employs built-in mathematical rule fallbacks if serialized models are missing during boot."
    ]
    for pt in points_ml1:
        p = tf_ml1.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)
        
    # Right Model: Linear Regression
    create_glass_card(slide8, Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_INDIGO)
    box_ml2 = slide8.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_ml2 = box_ml2.text_frame
    tf_ml2.word_wrap = True
    
    p_ml2h = tf_ml2.paragraphs[0]
    p_ml2h.text = "Upcoming Semester GPA Forecasting"
    p_ml2h.font.name = "Segoe UI"
    p_ml2h.font.size = Pt(18)
    p_ml2h.font.bold = True
    p_ml2h.font.color.rgb = ACCENT_INDIGO
    p_ml2h.space_after = Pt(14)
    
    points_ml2 = [
        "Algorithm: LinearRegression Model (Scikit-Learn).",
        "Input Indicators (3-dimensional): Class attendance ratios, assignment scores, and programming milestone rates.",
        "Mathematical Target: Forecasts predicted numerical GPA for the subsequent academic semester.",
        "Validation Tracking: Employs R-squared ($R^2$) scores and Mean Squared Error (MSE) metrics during retraining phases.",
        "Dynamic Re-Triggering: Instantly adjusts forecasts when students update tasks or grades change in real-time."
    ]
    for pt in points_ml2:
        p = tf_ml2.add_paragraph()
        p.text = "• " + pt
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # =========================================================================
    # SLIDE 9: PERSONALIZATION ENGINE
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide9, BG_COLOR)
    add_slide_header(slide9, "Personalization Heuristics & Adaptive Task Generation", "Core Algorithms")
    
    # Large detailed explanation card
    create_glass_card(slide9, Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8), border_color=ACCENT_SKY)
    box_pe = slide9.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.4))
    tf_pe = box_pe.text_frame
    tf_pe.word_wrap = True
    
    p_peh = tf_pe.paragraphs[0]
    p_peh.text = "How Personalization Works Behind the Scenes"
    p_peh.font.name = "Segoe UI"
    p_peh.font.size = Pt(18)
    p_peh.font.bold = True
    p_peh.font.color.rgb = ACCENT_SKY
    p_peh.space_after = Pt(14)
    
    pe_steps = [
        ("Step 1: Metric Checkpoints", "The system constantly scans individual records: Attendance, Assignment scoring averages, and subject marks."),
        ("Step 2: Objective Alignment", "Based on career goal selections, the system analyzes coding milestones. (e.g. For AIML goals: coding score <50 initiates core math/Python tasks; >75 triggers deep learning PyTorch deployment)."),
        ("Step 3: Weak-Subject Focus", "Checks subject score indexes. If scores in courses like Data Structures lie below 60%, the engine triggers targeted revision and past paper revision tasks."),
        ("Step 4: Risk Mitigation", "If the ML engine flags 'High Risk', critical counselor mentoring actions and 100% attendance targets are forced onto the weekly checklist dynamically."),
        ("Step 5: Interactive Syncing", "Regenerating tasks cleans current listings, computes modern scores, and populates optimized workflows restricted to 4 actionable tasks for readability.")
    ]
    
    for step_title, step_desc in pe_steps:
        p = tf_pe.add_paragraph()
        p.text = f"📍 {step_title} → "
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_INDIGO
        
        run = p.add_run()
        run.text = step_desc
        run.font.name = "Segoe UI"
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 10: CREDENTIALS & ROADMAP
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide10, BG_COLOR)
    add_slide_header(slide10, "Evaluation Access Tiers & Platform Roadmap", "Launch Plan")
    
    # Left Side: Test Credentials Table Card
    create_glass_card(slide10, Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_SKY)
    box_cr = slide10.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_cr = box_cr.text_frame
    tf_cr.word_wrap = True
    
    p_crh = tf_cr.paragraphs[0]
    p_crh.text = "Pre-seeded Evaluation Accounts"
    p_crh.font.name = "Segoe UI"
    p_crh.font.size = Pt(18)
    p_crh.font.bold = True
    p_crh.font.color.rgb = ACCENT_SKY
    p_crh.space_after = Pt(14)
    
    credentials = [
        ("Student Access Tier", "Login: student@edupilot.com\nDetails: Yaseen Ashu | CGPA 7.8 | 78% Attendance"),
        ("Faculty Access Tier", "Login: teacher@edupilot.com\nDetails: Dr. Sarah Connor | CS Department"),
        ("Administrator Access Tier", "Login: admin@edupilot.com\nDetails: Complete ML Control Panel Panel"),
        ("Global Password", "password123 (Applied across all test roles)")
    ]
    for role, details in credentials:
        p = tf_cr.add_paragraph()
        p.text = f"👤 {role}:\n"
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        
        run = p.add_run()
        run.text = details
        run.font.name = "Segoe UI"
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)
        
    # Right Side: Roadmap Timeline Card
    create_glass_card(slide10, Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8), border_color=ACCENT_INDIGO)
    box_rm = slide10.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(4.9), Inches(4.4))
    tf_rm = box_rm.text_frame
    tf_rm.word_wrap = True
    
    p_rmh = tf_rm.paragraphs[0]
    p_rmh.text = "Future Platform Roadmap"
    p_rmh.font.name = "Segoe UI"
    p_rmh.font.size = Pt(18)
    p_rmh.font.bold = True
    p_rmh.font.color.rgb = ACCENT_INDIGO
    p_rmh.space_after = Pt(14)
    
    roadmap_points = [
        ("Phase 1: LLM AI Mentors", "Deploy localized LLMs to offer detailed interactive feedback on specific weak subjects directly on the student dashboard."),
        ("Phase 2: Live LMS Synchronization", "Build API bridges to ingest real-time marks and attendance databases from platforms like Canvas, Moodle, or campus ERPs."),
        ("Phase 3: Recruiters Bridge", "Sync student placement readiness scores directly with recruiter job postings, matching candidates automatically to target placements.")
    ]
    for phase_name, phase_desc in roadmap_points:
        p = tf_rm.add_paragraph()
        p.text = f"🚀 {phase_name}:\n"
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        
        run = p.add_run()
        run.text = phase_desc
        run.font.name = "Segoe UI"
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    # 4. Save presentation to the parent Desktop directory (where the workspace folder resides)
    # The workspace folder is located at c:\Users\MOSHIN\OneDrive\Desktop\EduPilot
    # Desktop is the parent folder: c:\Users\MOSHIN\OneDrive\Desktop
    output_path = r"c:\Users\MOSHIN\OneDrive\Desktop\EduPilot_AI_Presentation.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation successfully created and saved to: {output_path}")
    except Exception as e:
        # Fallback to current directory if there is any folder access issue
        fallback_path = "EduPilot_AI_Presentation.pptx"
        prs.save(fallback_path)
        print(f"Presentation saved to fallback path due to: {e}\nLocation: {os.path.abspath(fallback_path)}")

if __name__ == "__main__":
    create_presentation()
